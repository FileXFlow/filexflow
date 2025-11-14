import os
import io
import tempfile
import zipfile

import pdfplumber
import fitz  # PyMuPDF
import img2pdf
from openpyxl import Workbook
from docx import Document

LIBREOFFICE_BIN = os.environ.get("LIBREOFFICE_BIN", "soffice")


# ---------- PDF → DOCX / XLSX / CSV / HTML / PPTX / JPG ZIP ----------

def pdf_to_docx(pdf_bytes: bytes) -> bytes:
    """
    Convierte un PDF en un DOCX simple (solo texto, sin formato complejo),
    usando pdfplumber + python-docx. NO usa LibreOffice.
    """
    doc = Document()

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        num_pages = len(pdf.pages)
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                for line in text.splitlines():
                    doc.add_paragraph(line)
            else:
                doc.add_paragraph("[Page without extractable text]")

            if i != num_pages:
                doc.add_page_break()

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def pdf_to_csv(pdf_bytes: bytes) -> bytes:
    """
    Extrae el texto de cada página línea por línea y lo vuelca como CSV 1 columna.
    """
    output = io.StringIO(newline="")
    import csv as _csv
    writer = _csv.writer(output)

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            for line in text.splitlines():
                writer.writerow([line])

    return output.getvalue().encode("utf-8")


def pdf_to_xlsx(pdf_bytes: bytes) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "PDF"

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        r = 1
        for i, page in enumerate(pdf.pages, start=1):
            ws.cell(r, 1, f"Page {i}")
            r += 1
            text = page.extract_text() or ""
            for line in text.splitlines():
                ws.cell(r, 1, line)
                r += 1
            r += 1

    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def pdf_to_html(pdf_bytes: bytes) -> bytes:
    """
    HTML muy simple: Page N + texto en <pre>. NO usa LibreOffice.
    """
    import html as _html

    parts: list[str] = ["<html><body>"]
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            parts.append(f"<h2>Page {i}</h2>")
            parts.append("<pre>")
            parts.append(_html.escape(text))
            parts.append("</pre>")
    parts.append("</body></html>")
    return "\n".join(parts).encode("utf-8")


def pdf_to_pptx(pdf_bytes: bytes) -> bytes:
    """
    Cada página del PDF como una diapositiva con la imagen de la página completa.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    for page in doc:
        slide = prs.slides.add_slide(blank)
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
            tf.write(img_bytes)
            tmp_path = tf.name

        slide.shapes.add_picture(
            tmp_path,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        os.unlink(tmp_path)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def pdf_to_jpg_zip(pdf_bytes: bytes) -> bytes:
    """
    Devuelve un ZIP con una imagen JPG por página.
    """
    zbio = io.BytesIO()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    with zipfile.ZipFile(zbio, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=150)
            jpg = pix.tobytes("jpg")
            zf.writestr(f"page-{i}.jpg", jpg)

    return zbio.getvalue()


# ---------- JPG → PDF ----------

def image_to_pdf(jpg_bytes: bytes) -> bytes:
    return img2pdf.convert(jpg_bytes)


# ---------- DOCX / XLSX / PPTX / HTML → PDF (aquí sí usamos LibreOffice) ----------

def office_to_pdf(file_bytes: bytes, ext: str) -> bytes:
    """
    Usa LibreOffice para convertir DOCX/XLSX/PPTX/HTML en PDF.
    """
    with tempfile.TemporaryDirectory() as tmp:
        in_path = os.path.join(tmp, f"in.{ext}")
        out_dir = tmp

        with open(in_path, "wb") as f:
            f.write(file_bytes)

        # soffice --headless --convert-to pdf in.ext --outdir tmp
        cmd = [
            LIBREOFFICE_BIN,
            "--headless",
            "--convert-to",
            "pdf",
            in_path,
            "--outdir",
            out_dir,
        ]
        # Si LibreOffice falla, lanzará excepción
        os.system(" ".join(cmd))

        out_path = os.path.join(out_dir, "in.pdf")
        if not os.path.exists(out_path):
            raise FileNotFoundError("LibreOffice did not produce in.pdf")

        with open(out_path, "rb") as f:
            return f.read()
